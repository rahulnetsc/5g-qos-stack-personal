"""Build the slide deck for the two-tier QoS-aware scheduler design.

Run with system python3 (python-pptx lives there, not in the uv venv):

    python3 deck/build_deck.py

The deck presents the design and what it buys against the proportional-fair
scheduler OpenAirInterface ships. Every number here is post-2026-08-07, i.e.
after the fidelity corrections recorded in design-docs/oai-phase1-review.md.
If a figure changes, change it here and rebuild -- do not edit the .pptx by
hand, or the deck and the repo will drift.

Colours are slots 1 and 2 of the dataviz reference categorical palette, used
in fixed order. That order is documented as pre-validated as a set (worst
adjacent CVD dE 9.1 in light mode); no JS runtime was available on this
machine to re-run the validator, so the fixed order is relied on rather than
re-derived.
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- design tokens ---------------------------------------------------------
INK = RGBColor(0x1A, 0x1A, 0x19)        # text-primary
INK_2 = RGBColor(0x5C, 0x5B, 0x54)      # text-secondary
INK_3 = RGBColor(0x8A, 0x89, 0x80)      # muted
RULE = RGBColor(0xDD, 0xDC, 0xD6)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF5, 0xF4, 0xF0)       # table header / callout fill
SERIES_1 = RGBColor(0x2A, 0x78, 0xD6)   # blue  -- PF
SERIES_2 = RGBColor(0xEB, 0x68, 0x34)   # orange -- TwoTier
GOOD = RGBColor(0x00, 0x83, 0x00)
BAD = RGBColor(0xE3, 0x49, 0x48)

FONT = "Calibri"
MONO = "Consolas"
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
BODY_W = W - 2 * MARGIN


# --- layout estimation -----------------------------------------------------
# python-pptx cannot measure text, and PowerPoint will happily render a
# paragraph straight out the bottom of its box and over whatever is below.
# These estimate rendered height so blocks can stack themselves; the
# per-point character widths are calibrated against LibreOffice's rendering
# of Carlito/Consolas at this deck's sizes, with headroom.
CW_PROP = 0.0069        # inches of advance per point of size, proportional
CW_MONO = 0.0077        # ditto, monospace
LINE_FUDGE = 1.16       # LibreOffice leads a little taller than nominal


def _para_lines(content, usable_in, size, cw=CW_PROP):
    per_line = max(8, int(usable_in / (cw * size)))
    return sum(max(1, -(-len(seg.replace("**", "")) // per_line))
               for seg in content.split("\n"))


def est_h(runs, usable_in, size=18, space_after=6, line=1.25, cw=CW_PROP):
    """Estimated rendered height, in inches, of what text() would draw."""
    total = 0.0
    for item in (runs if isinstance(runs, list) else [runs]):
        content, over = item if isinstance(item, tuple) else (item, {})
        sz = over.get("size", size)
        lines = _para_lines(content, usable_in, sz,
                            CW_MONO if over.get("font") == MONO else cw)
        total += lines * sz * over.get("line", line) * LINE_FUDGE / 72.0
        total += over.get("space_after", space_after) / 72.0
    return total


def text(slide, x, y, w, h, runs, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space_after=6, line=1.25, font=FONT):
    """runs: a string, or a list of (text, {overrides}) tuples per paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    items = runs if isinstance(runs, list) else [runs]
    for i, item in enumerate(items):
        content, over = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.space_after = Pt(over.get("space_after", space_after))
        p.line_spacing = over.get("line", line)
        # inline bold via ** markers
        for j, seg in enumerate(content.split("**")):
            if not seg:
                continue
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(over.get("size", size))
            r.font.color.rgb = over.get("color", color)
            r.font.name = over.get("font", font)
            r.font.bold = (j % 2 == 1) or over.get("bold", bold)
    return box


def rule(slide, y, x=MARGIN, w=None, color=RULE, h=Pt(1.25)):
    w = w or BODY_W
    s = slide.shapes.add_shape(1, x, y, w, h)   # 1 = rectangle
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def slide_base(prs, title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    y = Inches(0.55)
    if kicker:
        text(s, MARGIN, y, BODY_W, Inches(0.3), kicker.upper(),
             size=11, color=INK_3, bold=True)
        y += Inches(0.34)
    text(s, MARGIN, y, BODY_W, Inches(0.7), title, size=30, bold=True)
    rule(s, Inches(1.62))
    return s


def table(slide, x, y, w, rows, col_w=None, header=True, size=14,
          emphasis=None, numeric=True, row_h=0.34):
    """rows: list of lists of str. emphasis: {(r,c): RGBColor}.

    numeric=True right-aligns every column but the first (the default, for
    figures); numeric=False left-aligns throughout, for tables of prose.
    """
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, x, y, w, Inches(row_h) * nr)
    tbl = shape.table
    if col_w:
        total = sum(col_w)
        for i, frac in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * frac / total))
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if (header and r == 0) else SURFACE
            cell.margin_left = Inches(0.09)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            p.alignment = (PP_ALIGN.RIGHT if (numeric and c > 0)
                           else PP_ALIGN.LEFT)
            r_ = p.add_run()
            r_.text = val
            r_.font.size = Pt(size)
            r_.font.name = FONT
            r_.font.bold = (header and r == 0)
            r_.font.color.rgb = (emphasis or {}).get((r, c), INK if r == 0 else INK_2)
    return shape


def block(slide, x, y, w, runs, gap=0.18, **kw):
    """Place a self-sized text block. Returns the y just below it."""
    h = est_h(runs, w / 914400.0 - 0.08, **{k: v for k, v in kw.items()
                                            if k in ("size", "space_after",
                                                     "line")})
    text(slide, x, y, w, Inches(h), runs, **kw)
    return y + Inches(h + gap)


def callout(slide, x, y, w, body, accent=SERIES_2, size=16, gap=0.2):
    """Self-sized tinted panel with an accent bar. Returns the y below it."""
    h = Inches(est_h(body, w / 914400.0 - 0.58, size=size) + 0.40)
    bg = slide.shapes.add_shape(1, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BAND
    bg.line.fill.background()
    bg.shadow.inherit = False
    bar = slide.shapes.add_shape(1, x, y, Pt(4), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.5),
         h - Inches(0.3), body, size=size, color=INK)
    return y + h + Inches(gap)


def math(slide, x, y, w, lines, size=14, line=1.3, gap=0.18):
    """A block of set-in-monospace formulae. lines: list of (expr, gloss)."""
    runs = []
    for expr, gloss in lines:
        runs.append((expr, {"font": MONO, "size": size, "color": INK,
                            "space_after": 0 if gloss else 7, "line": line}))
        if gloss:
            runs.append((gloss, {"size": size - 3, "color": INK_3,
                                 "space_after": 7, "line": 1.1}))
    return block(slide, x, y, w, runs, gap=gap)


# ===========================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = W, H


# Column geometry used by the two-column slides.
COL_L, COL_R = MARGIN, Inches(7.1)
COL_W = Inches(5.3)
BODY_TOP = Inches(1.88)

# --- 1. title --------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
text(s, MARGIN, Inches(2.3), BODY_W, Inches(1.6),
     "Design of a Two-Tier QoS-Aware Scheduler\nfor Private 5G",
     size=36, bold=True, line=1.18)
rule(s, Inches(4.0), w=Inches(2.2), color=SERIES_2, h=Pt(3))
block(s, MARGIN, Inches(4.3), BODY_W, [
    ("A convex rate planner on a ~1 s horizon, over a per-slot "
     "drift-plus-penalty tracker.", {"size": 19, "color": INK_2}),
    ("Motivation, design, and simulation results against the scheduler "
     "OpenAirInterface ships — with what remains to land it on hardware.",
     {"size": 19, "color": INK_2}),
])

# --- 2. motivation ---------------------------------------------------------
s = slide_base(prs, "The factory floor asks for three different things",
               "motivation")
y = block(s, MARGIN, Inches(1.95), BODY_W,
          "One private 5G cell, one carrier, and traffic classes that differ "
          "in kind rather than in volume:", size=18, color=INK_2, gap=0.28)
table(s, MARGIN, y, BODY_W, [
    ["Class", "Example", "What it needs", "How you know it failed"],
    ["GBR", "Robot camera and LIDAR uplink", "A sustained rate floor",
     "Dropped frames, blind robot"],
    ["Delay", "Motion control, teleoperation", "Bytes before a deadline",
     "A late command — a safety event"],
    ["Best-effort", "Firmware push, log upload", "Whatever is left",
     "Slower; nobody notices"],
], col_w=[1.1, 2.5, 2.0, 2.6], size=15, numeric=False)
y += Inches(0.34 * 4 + 0.3)
y = callout(s, MARGIN, y, BODY_W,
            "OpenAirInterface's NR MAC schedules with a **proportional-fair "
            "metric over per-UE grants**. PF has no term for a rate contract, "
            "no term for a deadline, and no configured-grant path for "
            "periodic traffic. It is a good answer to a question the factory "
            "is not asking.", gap=0.3)
block(s, MARGIN, y, BODY_W,
      "A robot whose video is switched off entirely is a failure in a way a "
      "uniformly degraded fleet is not — a distinction a throughput-fair "
      "metric cannot see, let alone act on.", size=16, color=INK_3)

# --- 3. state of the art ---------------------------------------------------
s = slide_base(prs, "What the prior art told us to do — and not to",
               "informing the design")
yl = block(s, COL_L, BODY_TOP, Inches(5.9), [
    ("Adopted", {"size": 18, "bold": True, "color": GOOD}),
    ("**Two timescales.** A slow declarative program over a fast per-slot "
     "controller — the pattern behind NVS (Kokku 2012) and the O-RAN RIC "
     "split. We instantiate it; we do not claim it.",
     {"size": 15, "color": INK_2}),
    ("**Network utility maximisation** (Kelly 1998) for the slow tier — an "
     "objective a plant engineer can read and audit.",
     {"size": 15, "color": INK_2}),
    ("**Drift-plus-penalty** (Neely 2010) for the fast tier — it solves the "
     "per-slot problem, and needs a target, which Tier-1 supplies.",
     {"size": 15, "color": INK_2}),
    ("**Configured grants** dominate latency for time-critical industrial "
     "traffic (Larrañaga 2023, ns-3 / 5G-LENA).",
     {"size": 15, "color": INK_2}),
])
yr = block(s, COL_R, BODY_TOP, COL_W, [
    ("Ruled out, and why", {"size": 18, "bold": True, "color": BAD}),
    ("**Tuned PF variants** — M-LWDF, Exp-rule. Excellent on delay; on "
     "rate they saturate, because a multiplicative deficit weight equalises "
     "a weighted 1/R_avg and cannot reach an arbitrary unequal target. "
     "Monghal (2008) measured exactly that.", {"size": 15, "color": INK_2}),
    ("**Learned schedulers.** Kela (2024), arguing from inside that "
     "literature, finds deployable deep schedulers do not yet exist — "
     "real-time cost, 3GPP compliance, and poor generalisation across "
     "bandwidth and traffic mix.", {"size": 15, "color": INK_2}),
])
callout(s, MARGIN, max(yl, yr) + Inches(0.15), BODY_W,
        "Two commitments follow: an **integrator, not a metric "
        "multiplier** in the fast tier, so an arbitrary target vector is "
        "reachable; and a **convex program with one auditable fairness "
        "dial** in the slow tier.")

# --- 4. Tier-1 -------------------------------------------------------------
s = slide_base(prs, "Tier-1: what rate should each flow target?",
               "the design  ·  every ~1 s")
y = block(s, COL_L, Inches(1.82), Inches(6.3),
          "Rates r_i ≥ 0, slacks s ≥ 0, one feasible set F(floor):",
          size=15, color=INK_2, gap=0.14)
y = math(s, COL_L, y, Inches(6.3), [
    ("Σ_{i∈d} r_i / SE_i  ≤  C_d",
     "capacity in PRB-symbols, special slot included"),
    ("r_i + s_i  ≥  G_i",
     "GBR contract — soft, so overload stays feasible"),
    ("r_i  ≥  floor_i",
     "max-min floor — hard"),
    ("Σ_{i∈σ} r_i/SE_i + s_σ  ≥  φ_σ·C_d",
     "slice share — soft, demand-capped, work-conserving"),
], gap=0.16)
y = block(s, COL_L, y, Inches(6.3), "Solved in three stages, in that order:",
          size=15, color=INK_2, gap=0.14)
math(s, COL_L, y, Inches(6.3), [
    ("A.  t* = max t  s.t.  r_i ≥ t·Ĝ_i  ∀ GBR",
     "then  floor_i = α·t*·Ĝ_i"),
    ("B1. S* = min Σ p_i·s_i + p_σ Σ SE_σ·s_σ", ""),
    ("B2.      max Σ w_c·log(r_i+ε)  s.t. Σ p_i·s_i ≤ S*", ""),
])
y = callout(s, Inches(7.6), Inches(1.85), Inches(4.85),
            "**Order the objectives; do not weight them.** Folded into "
            "one, the two terms sit 3.3×10⁷ apart — two solvers returned "
            "optimal_inaccurate and disagreed 3.6× on one flow's rate. Split "
            "into B1 then B2, both match the analytic optimum to 100 bps.",
            size=15)
callout(s, Inches(7.6), y, Inches(4.85),
        "**t * is the largest fraction of contract every GBR flow can "
        "hold at once.** t* = 1 means the set is feasible, the floor binds "
        "nothing, and stage A is free — which is why it is on by default. "
        "α trades the guarantee against throughput.",
        accent=GOOD, size=15)

# --- 5. Tier-2 and the MAC -------------------------------------------------
s = slide_base(prs, "Tier-2: hit that target, slot by slot",
               "the design  ·  every slot")
y = block(s, COL_L, Inches(1.82), Inches(6.3),
          "Every flow carries a virtual queue — an accumulator in bits, "
          "not a buffer of data:", size=15, color=INK_2, gap=0.14)
y = math(s, COL_L, y, Inches(6.3), [
    ("Q_i += r_i·Δ ;  Q_i ← min(Q_i, ceil_i)",
     "grow by the Tier-1 target, then clamp"),
    ("serve ;  Q_i ← max(0, Q_i − delivered_i)", ""),
    ("ceil_i = max(0, min(r_i·W, A_i^W) − D_i^W)",
     "the debt the flow legitimately accrued over the window W"),
], gap=0.16)
y = block(s, COL_L, y, Inches(6.3), "Ranking, and then the transport block:",
          size=15, color=INK_2, gap=0.14)
math(s, COL_L, y, Inches(6.3), [
    ("Q̃_i = Q_i + w_d·(HoL_i/PDB_i)^κ · max_j Q_j",
     "deadline urgency, scaled by the system-wide max queue"),
    ("M_u = ( Σ_{i∈u} Q̃_i ) · SE_u",
     "rank UEs — opportunistic and rate-tracking at once"),
    ("fill the TB by TS 38.321 LCP: priority, then Q̃_i", ""),
])
y = callout(s, Inches(7.6), Inches(1.85), Inches(4.85),
            "Q_i is a **non-saturating integrator**, so it can reach an "
            "arbitrary, unequal target vector — precisely what a "
            "multiplicative urgency metric cannot do at any weight.",
            size=15)
y = callout(s, Inches(7.6), y, Inches(4.85),
            "Clamping to a **windowed arrival count**, not instantaneous "
            "backlog: a bursty video flow empties between frames, and "
            "clamping there erases its legitimate debt.",
            accent=INK_3, size=15)
callout(s, Inches(7.6), y, Inches(4.85),
        "**Configured grants** carry the periodic flows — standing "
        "allocations in priority tiers, self-gating where they would starve "
        "the dynamic pool.", accent=GOOD, size=15)

# --- 6. the simulator ------------------------------------------------------
s = slide_base(prs, "The simulator: what it models, and what it does not",
               "evaluation")
yl = block(s, COL_L, BODY_TOP, Inches(5.9), [
    ("Modelled", {"size": 17, "bold": True, "color": GOOD}),
    ("• The slot-by-slot PRB grid, with the TDD special slot",
     {"size": 15, "color": INK_2}),
    ("• The PDCCH/CCE budget, with DCI aggregation levels",
     {"size": 15, "color": INK_2}),
    ("• The uplink BSR round-trip — delay plus loss",
     {"size": 15, "color": INK_2}),
    ("• CQI delay and loss; BLER from MCS/SNR mismatch",
     {"size": 15, "color": INK_2}),
    ("• Per-UE AR(1) SNR through an MCS staircase",
     {"size": 15, "color": INK_2}),
    ("• Per-flow buffers, HoL timestamps, PDB expiry",
     {"size": 15, "color": INK_2}),
])
yr = block(s, COL_R, BODY_TOP, COL_W, [
    ("Not modelled", {"size": 17, "bold": True, "color": INK_3}),
    ("• LDPC and I/Q; HARQ as a state machine — a fixed delay and a "
     "BLER discount instead", {"size": 15, "color": INK_2}),
    ("• Per-packet RLC segmentation, MU-MIMO", {"size": 15, "color": INK_2}),
    ("• Mobility, UE churn, inter-cell interference",
     {"size": 15, "color": INK_2}),
    ("So the results are **comparative, not absolute** — and what we omit "
     "would widen the configured-grant margin, not narrow it.",
     {"size": 15, "color": INK_3}),
])
y = block(s, MARGIN, max(yl, yr) + Inches(0.05), BODY_W,
          "Three workloads, each sized so a different bottleneck binds:",
          size=15, color=INK_2, gap=0.24)
table(s, MARGIN, y, BODY_W, [
    ["Scenario", "Carrier / TDD", "Workload", "What binds"],
    ["factory_robots", "40 MHz, μ=2, DSUUU", "10 robots, 24 flows, 14–24 dB",
     "Uplink data channel"],
    ["sensor_dense", "30 MHz, μ=1, DSUUU", "30 periodic sensors, 15 ms PDB",
     "PDCCH / CCE budget"],
    ["latency_bound", "40 MHz, μ=1, DDDSU", "8 × 5 Mbps + 80 Mbps bulk",
     "Deadlines under saturation"],
], col_w=[1.5, 1.9, 2.6, 1.9], size=14, numeric=False, row_h=0.32)

# --- 7. result: control channel --------------------------------------------
s = slide_base(prs, "Result 1: a capability PF does not have",
               "sensor_dense  ·  the control channel binds")
y = block(s, MARGIN, Inches(1.95), BODY_W,
          "30 dense periodic uplink sensors with a 15 ms deadline. The "
          "DCI/CCE budget runs out before the data channel does.",
          size=17, color=INK_2, gap=0.3)
table(s, MARGIN, y, Inches(8.6), [
    ["Scheduler", "Total", "On time", "Min delivery", "Worst p99"],
    ["Round Robin", "6.9 M", "0/30", "71%", "15.0 ms"],
    ["Proportional Fair", "9.2 M", "2/30", "85%", "15.0 ms"],
    ["Two-tier", "9.6 M", "30/30", "100%", "5.0 ms"],
], col_w=[2.6, 1.3, 1.3, 1.7, 1.6], size=16,
   emphasis={(3, 1): GOOD, (3, 2): GOOD, (3, 3): GOOD, (3, 4): GOOD})
y += Inches(0.34 * 4 + 0.3)
y = callout(s, MARGIN, y, BODY_W,
            "Each periodic flow gets a standing allocation that costs **zero "
            "PDCCH per slot and needs no BSR round trip**. PF is throttled by "
            "both budgets at once and has no mechanism to bypass either — so "
            "no amount of PF tuning closes this.", accent=GOOD, gap=0.3)
block(s, MARGIN, y, BODY_W,
      "Larrañaga et al. (JNCA 2023) report the same dominance of configured "
      "grants for time-critical industrial traffic from ns-3 / 5G-LENA — an "
      "independent simulator sharing no code with ours.",
      size=15, color=INK_3)

# --- 8. result: deadlines --------------------------------------------------
s = slide_base(prs, "Result 2: deadline-blindness is silent",
               "latency_bound  ·  deadlines under saturation")
y = block(s, MARGIN, Inches(1.95), BODY_W,
          "Eight interactive 5 Mbps streams with a 12 ms deadline, sharing a "
          "saturated downlink with 80 Mbps of bulk.",
          size=17, color=INK_2, gap=0.3)
table(s, MARGIN, y, Inches(9.4), [
    ["Scheduler", "On time", "Mean delivery", "Worst p99", "Bulk carried"],
    ["Round Robin", "3/8", "84%", "12.0 ms", "22.8 M"],
    ["Proportional Fair", "5/8", "86%", "12.0 ms", "24.6 M"],
    ["Two-tier", "8/8", "100%", "4.5 ms", "13.2 M"],
], col_w=[2.6, 1.2, 1.6, 1.4, 1.5], size=16,
   emphasis={(3, 1): GOOD, (3, 2): GOOD, (3, 3): GOOD})
y += Inches(0.34 * 4 + 0.3)
y = callout(s, MARGIN, y, BODY_W,
            "PF's mean delivery on the control flows is **86% — which reads "
            "healthy on a dashboard.** The missing 14% are aged-out packets: "
            "the late motion commands, the safety-relevant ones. Mean "
            "delivery is not a safety metric.", accent=BAD, gap=0.3)
block(s, MARGIN, y, BODY_W,
      "The design pays for this explicitly and visibly — bulk throughput "
      "drops 24.6 → 13.2 Mbps to clear every deadline. That is the trade the "
      "QoS profile asked for.", size=15, color=INK_3)

# --- 9. result: GBR under load ---------------------------------------------
s = slide_base(prs, "Result 3: rate contracts, swept across offered load",
               "factory_robots  ·  the uplink data channel binds")
block(s, COL_L, Inches(1.88), Inches(6.0),
      "Ten robots, uplink-heavy. Capacity is fixed by spectrum, so the swept "
      "axis is offered load.", size=17, color=INK_2)

chart_data = CategoryChartData()
chart_data.categories = ["1.00x", "0.67x", "0.50x", "0.33x"]
chart_data.add_series("Proportional Fair", (1, 21, 52, 83))
chart_data.add_series("Two-tier", (44, 67, 81, 83))
chart_shape = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, Inches(2.55),
    Inches(5.9), Inches(3.4), chart_data)
ch = chart_shape.chart
ch.has_title = True
ch.chart_title.text_frame.text = "Worst-served GBR flow (% of contract)"
tp = ch.chart_title.text_frame.paragraphs[0]
tp.runs[0].font.size = Pt(14)
tp.runs[0].font.bold = False
tp.runs[0].font.color.rgb = INK_2
tp.runs[0].font.name = FONT
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(12)
ch.legend.font.name = FONT
ch.legend.font.color.rgb = INK_2
ch.plots[0].gap_width = 60
ch.plots[0].series[0].format.fill.solid()
ch.plots[0].series[0].format.fill.fore_color.rgb = SERIES_1
ch.plots[0].series[1].format.fill.solid()
ch.plots[0].series[1].format.fill.fore_color.rgb = SERIES_2
va = ch.value_axis
va.has_major_gridlines = True
va.maximum_scale = 100
va.tick_labels.font.size = Pt(11)
va.tick_labels.font.color.rgb = INK_3
ca = ch.category_axis
ca.tick_labels.font.size = Pt(11)
ca.tick_labels.font.color.rgb = INK_3

y = block(s, COL_R, Inches(2.35), COL_W,
          "Contracts met (≥95% of the rate floor)", size=14, color=INK_2,
          gap=0.12)
table(s, COL_R, y, COL_W, [
    ["Load", "PF", "Two-tier"],
    ["1.00x  (as shipped)", "1/10", "0/10"],
    ["0.67x", "6/10", "1/10"],
    ["0.50x", "8/10", "10/10"],
    ["0.33x", "10/10", "10/10"],
], col_w=[2.6, 1.2, 1.4], size=15,
   emphasis={(2, 1): GOOD, (3, 2): GOOD})
y += Inches(0.34 * 5 + 0.3)
callout(s, COL_R, y, COL_W,
        "**At 0.50× we honour 10/10 contracts to PF's 8/10.** Deeper in "
        "overload the two disagree about the metric, not the outcome: at "
        "0.67× every robot holds ≥67% of its rate where PF's worst gets 21%.",
        size=15)
block(s, MARGIN, Inches(6.75), BODY_W,
      "A threshold count rewards concentrating the shortfall; the max-min "
      "floor spreads it, and α picks between them. This deep into overload "
      "the operator's real lever is admission control.",
      size=15, color=INK_3)

# --- 10. what it adds ------------------------------------------------------
s = slide_base(prs, "What the design adds — and what it costs", "summary")
yl = block(s, COL_L, BODY_TOP, Inches(5.9), [
    ("Added over the PF baseline", {"size": 17, "bold": True, "color": GOOD}),
    ("• A **rate contract** that is stated and honoured — 10/10 against "
     "8/10 at moderate load; worst-served 44% against 1% in overload.",
     {"size": 15, "color": INK_2}),
    ("• A deadline that is **tracked, not hoped for** — 8/8 against 5/8, "
     "p99 4.5 against 12.0 ms.", {"size": 15, "color": INK_2}),
    ("• A **control-channel path** for dense periodic traffic — 30/30 "
     "against 2/30 where the CCE budget binds.", {"size": 15, "color": INK_2}),
    ("• A **single fairness dial** that self-disables when the GBR set is "
     "feasible, plus soft per-slice floors.", {"size": 15, "color": INK_2}),
])
yr = block(s, COL_R, BODY_TOP, COL_W, [
    ("Paid for with", {"size": 17, "bold": True, "color": SERIES_2}),
    ("• About **4% of aggregate throughput** at the max-min floor in deep "
     "overload — the price of holding the worst-served flow above zero.",
     {"size": 15, "color": INK_2}),
    ("• **One convex solve per second** — CVXPY/Clarabel here, GLPK simplex "
     "with successive convex approximation in the OAI port.",
     {"size": 15, "color": INK_2}),
    ("• Per-5QI priority and prioritised-bit-rate configuration that has to "
     "actually be set.", {"size": 15, "color": INK_2}),
])
callout(s, MARGIN, max(yl, yr) + Inches(0.1), BODY_W,
        "It is also **markedly less sensitive to control-plane fidelity.** "
        "Sweeping the uplink BSR round-trip from 0 to 16 slots, PF's "
        "contract count breaks 8/10 → 5/10 and its worst-served flow falls "
        "66% → 30%; the two-tier design moves 10/10 → 8/10 and 80% → 77%, "
        "losing 1% of throughput against PF's 10%. Configured-grant flows "
        "are invariant to BSR delay by construction.", accent=GOOD, size=15)
# --- 11. status and next ---------------------------------------------------
s = slide_base(prs, "Where this stands, and what remains", "status")
C1, C2, C3 = MARGIN, Inches(4.85), Inches(8.8)
CW3 = Inches(3.6)
ys = []
ys.append(block(s, C1, BODY_TOP, CW3, [
    ("Built and tested", {"size": 17, "bold": True, "color": GOOD}),
    ("• Scheduler library: Tier-1, Tier-2, configured grants, per-UE grants "
     "with a TS 38.321 multiplexer, per-slice floors",
     {"size": 15, "color": INK_2}),
    ("• A symbol-accurate NR simulator with PDCCH, BSR and CQI models",
     {"size": 15, "color": INK_2}),
    ("• Three scenario studies, 70 regression tests, one script per result "
     "table", {"size": 15, "color": INK_2}),
]))
ys.append(block(s, C2, BODY_TOP, CW3, [
    ("In progress", {"size": 17, "bold": True, "color": SERIES_2}),
    ("• OAI integration in a separate repo — Tier-1 as GLPK simplex with "
     "successive convex approximation, Tier-2 inside the NR MAC",
     {"size": 15, "color": INK_2}),
    ("• Porting the max-min stage and the per-5QI priority mapping onto that "
     "branch", {"size": 15, "color": INK_2}),
]))
ys.append(block(s, C3, BODY_TOP, CW3, [
    ("To test on hardware", {"size": 17, "bold": True, "color": INK_3}),
    ("• End-to-end runs on OAI with COTS UEs", {"size": 15, "color": INK_2}),
    ("• Tier-1 solve time inside a real 1 s budget",
     {"size": 15, "color": INK_2}),
    ("• What the simulator omits: HARQ PRB cost, uplink grant timing, "
     "mobility, UE churn", {"size": 15, "color": INK_2}),
    ("• Multi-cell and inter-cell interference", {"size": 15, "color": INK_2}),
]))
callout(s, MARGIN, max(ys) + Inches(0.15), BODY_W,
        "The design is settled and the simulation case is made: against the "
        "PF scheduler OAI ships, it adds rate contracts, deadline tracking "
        "and a control-channel path, at a cost we can state and bound. "
        "**The open work is landing it on real hardware** — everything above "
        "is comparative, in a simulator, and silicon is the next step.")

prs.save("deck/two-tier-scheduler.pptx")
print(f"wrote deck/two-tier-scheduler.pptx  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
